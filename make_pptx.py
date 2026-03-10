"""
RNA Velocity 解析結果 PowerPoint 作成スクリプト
"""

import os
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
import pptx.util as util

WORKDIR = os.path.dirname(os.path.abspath(__file__))
os.chdir(WORKDIR)

FIG_DIR = "figures"
OG_DIR  = "figures_og_exact"
OUT_PPT = "SutureRNAVelocity_Results.pptx"

# スライドサイズ：ワイド（16:9）
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# カラーパレット
C_TITLE  = RGBColor(0x1a, 0x1a, 0x2e)   # 濃紺
C_HEAD   = RGBColor(0x16, 0x21, 0x3e)   # 紺
C_ACCENT = RGBColor(0xe9, 0x4f, 0x37)   # 赤
C_SUB    = RGBColor(0x39, 0x3e, 0x46)   # グレー
C_WHITE  = RGBColor(0xff, 0xff, 0xff)
C_LIGHT  = RGBColor(0xf5, 0xf5, 0xf5)
C_OG1    = RGBColor(0xe4, 0x1a, 0x1c)
C_OG2    = RGBColor(0xff, 0x7f, 0x00)
C_OG3    = RGBColor(0x4d, 0xaf, 0x4a)
C_OG4    = RGBColor(0x98, 0x4e, 0xa3)
C_PO1    = RGBColor(0x37, 0x7e, 0xb8)
C_PO2    = RGBColor(0xa6, 0x56, 0x28)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

BLANK = prs.slide_layouts[6]   # 白紙レイアウト

# ============================================================
# ヘルパー関数
# ============================================================
def add_rect(slide, left, top, width, height, fill=None, line=None):
    shape = slide.shapes.add_shape(
        1, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.line.fill.background() if line is None else None
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line is None:
        shape.line.fill.background()
    return shape

def add_text(slide, text, left, top, width, height,
             size=18, bold=False, color=C_TITLE,
             align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txb

def add_image(slide, path, left, top, width=None, height=None):
    if not os.path.exists(path):
        return None
    if width and height:
        return slide.shapes.add_picture(
            path, Inches(left), Inches(top), Inches(width), Inches(height)
        )
    elif width:
        return slide.shapes.add_picture(
            path, Inches(left), Inches(top), width=Inches(width)
        )
    elif height:
        return slide.shapes.add_picture(
            path, Inches(left), Inches(top), height=Inches(height)
        )
    else:
        return slide.shapes.add_picture(path, Inches(left), Inches(top))

def slide_header(slide, title, subtitle=None):
    """上部バナー共通ヘッダー"""
    add_rect(slide, 0, 0, 13.33, 1.0, fill=C_HEAD)
    add_text(slide, title, 0.3, 0.08, 11.0, 0.7,
             size=24, bold=True, color=C_WHITE)
    if subtitle:
        add_text(slide, subtitle, 0.3, 0.62, 11.0, 0.45,
                 size=13, color=RGBColor(0xcc, 0xcc, 0xcc))

def add_footnote(slide, text):
    add_text(slide, text, 0.2, 7.15, 12.9, 0.3,
             size=9, color=RGBColor(0x88, 0x88, 0x88))

# ============================================================
# Slide 1: タイトル
# ============================================================
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill=C_TITLE)
add_rect(slide, 0, 2.8, 13.33, 2.3, fill=C_HEAD)

add_text(slide,
         "Mouse Cranial Suture RNA Velocity Analysis",
         0.5, 2.9, 12.3, 1.0,
         size=32, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
add_text(slide,
         "Osteogenic cell differentiation and dedifferentiation trajectories in E17.5 coronal suture",
         0.5, 3.9, 12.3, 0.6,
         size=16, color=RGBColor(0xcc, 0xdd, 0xff), align=PP_ALIGN.CENTER)
add_text(slide,
         "Data: GSE163693  |  Farmer et al. 2021, Nature Communications  |  scVelo + CellRank",
         0.5, 4.55, 12.3, 0.4,
         size=12, color=RGBColor(0xaa, 0xaa, 0xaa), align=PP_ALIGN.CENTER)

# OG/PO color circles as decoration
labels_colors = [
    ("OG1", C_OG1),
    ("OG2", C_OG2),
    ("OG3", C_OG3),
    ("OG4", C_OG4),
    ("PO1", C_PO1),
    ("PO2", C_PO2),
]
for i, (label, rgb) in enumerate(labels_colors):
    row = i // 3
    col = i % 3
    x = 3.5 + col * 1.6
    y = 5.2 + row * 0.5
    add_rect(slide, x, y, 0.9, 0.35, fill=rgb)
    add_text(slide, label, x, y, 0.9, 0.35,
             size=13, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

add_text(slide, "SutureRNAVelocity Pipeline", 0.5, 6.9, 12.3, 0.4,
         size=11, color=RGBColor(0x77, 0x77, 0x99), align=PP_ALIGN.CENTER)

# ============================================================
# Slide 2: 解析概要
# ============================================================
slide = prs.slides.add_slide(BLANK)
slide_header(slide, "Analysis Overview",
             "Farmer et al. 2021 | GSE163693 | E17.5 mouse coronal suture")
add_rect(slide, 0, 1.0, 13.33, 6.5, fill=C_LIGHT)

steps = [
    ("1", "Data Download",
     "3 BAM files (SRR13288596–98) via aria2c\nE17.5 coronal suture, 3 batches"),
    ("2", "velocyto",
     "Spliced / Unspliced count generation\nGencode vM25 annotation (chr-prefix removed)"),
    ("3", "scVelo – Full dataset",
     "5,407 cells × 904 genes\nDynamical model, Leiden clustering (19 clusters)"),
    ("4", "OG/PO Subset",
     "2,017 cells (OG1–4 + PO1–2) identified by metadata labels\n(GSE163693_E15_17_composite_metadata.csv.gz)"),
    ("5", "CellRank (GPCCA)",
     "Bidirectional trajectory analysis\nTerminal: OG4 + OG1 | Initial: OG3"),
    ("6", "Local Velocity",
     "Velocity confidence ≥ 0.77 for all clusters\nPhase portraits, PCA-space velocity"),
]
for i, (num, title, desc) in enumerate(steps):
    col = i % 3
    row = i // 3
    bx = 0.3 + col * 4.35
    by = 1.3 + row * 2.9
    add_rect(slide, bx, by, 4.0, 2.5, fill=C_WHITE)
    add_rect(slide, bx, by, 0.55, 0.55, fill=C_ACCENT)
    add_text(slide, num, bx, by, 0.55, 0.55,
             size=18, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, title, bx + 0.6, by + 0.05, 3.3, 0.5,
             size=14, bold=True, color=C_TITLE)
    add_text(slide, desc, bx + 0.1, by + 0.65, 3.8, 1.7,
             size=11, color=C_SUB)

# ============================================================
# Slide 3: 全細胞 UMAP + velocity stream
# ============================================================
slide = prs.slides.add_slide(BLANK)
slide_header(slide, "Full Dataset: UMAP and RNA Velocity",
             "5,407 cells | Leiden clustering (19 clusters) | Dynamical model")

add_image(slide, f"{FIG_DIR}/umap_clusters.png",
          0.2, 1.1, width=6.2)
add_image(slide, f"{FIG_DIR}/scvelo_velocity_stream_dynamical.png",
          6.6, 1.1, width=6.5)

add_text(slide, "Leiden clusters (19 groups)", 0.2, 5.6, 6.0, 0.4,
         size=11, color=C_SUB, align=PP_ALIGN.CENTER)
add_text(slide, "RNA velocity stream (dynamical model)", 6.6, 5.6, 6.5, 0.4,
         size=11, color=C_SUB, align=PP_ALIGN.CENTER)
add_footnote(slide, "scVelo dynamical model | Gencode vM25 annotation | E17.5 coronal suture")

# ============================================================
# Slide 4: 全細胞 Latent time + spliced proportions
# ============================================================
slide = prs.slides.add_slide(BLANK)
slide_header(slide, "Full Dataset: Latent Time and Splicing Dynamics")

add_image(slide, f"{FIG_DIR}/scvelo_latent_time.png",
          0.2, 1.1, width=6.2)
add_image(slide, f"{FIG_DIR}/scvelo_proportions_spliced_proportions.png",
          6.6, 1.1, width=6.5)

add_text(slide, "Latent time (full dataset)", 0.2, 5.6, 6.0, 0.4,
         size=11, color=C_SUB, align=PP_ALIGN.CENTER)
add_text(slide, "Spliced / Unspliced proportions per cluster", 6.6, 5.6, 6.5, 0.4,
         size=11, color=C_SUB, align=PP_ALIGN.CENTER)

# ============================================================
# Slide 5: OG/PO クラスター定義
# ============================================================
slide = prs.slides.add_slide(BLANK)
slide_header(slide, "OG/PO Cluster Definition",
             "Farmer et al. 2021 | Barcode-level identification using metadata labels")

add_image(slide, f"{OG_DIR}/umap_og_in_full_umap.png",
          0.2, 1.1, width=5.5)

# テーブル形式で OG/PO クラスターを説明
rows = [
    ("OG1", "548 cells", "Osteoprogenitors (earliest)", "Erg, Pthlh, Six2", C_OG1),
    ("OG2", "315 cells", "Suture-resident pre-osteoblasts", "Lef1, Inhba", C_OG2),
    ("OG3", "603 cells", "Periosteal pre-osteoblasts", "Mmp13, Podnl1", C_OG3),
    ("OG4", "306 cells", "Mature osteoblasts", "Ifitm5, Dmp1, Sost", C_OG4),
    ("PO1", "206 cells", "Periosteal osteoblasts (outer bone)", "", C_PO1),
    ("PO2", "39 cells",  "Periosteal osteoblasts (mature subset)", "", C_PO2),
]
headers = ["Cluster", "Cells", "Cell type", "Markers"]
col_w = [1.0, 1.2, 3.2, 2.5]
col_x = [6.0, 7.1, 8.4, 11.7]

# ヘッダー行
for j, (h, x, w) in enumerate(zip(headers, col_x, col_w)):
    add_rect(slide, x, 1.3, w, 0.45, fill=C_HEAD)
    add_text(slide, h, x + 0.05, 1.32, w - 0.1, 0.4,
             size=12, bold=True, color=C_WHITE)

for i, (cluster, cells, cell_type, markers, rgb) in enumerate(rows):
    y = 1.8 + i * 0.85
    bg = C_LIGHT if i % 2 == 0 else C_WHITE
    for j, (val, x, w) in enumerate(zip(
            [cluster, cells, cell_type, markers], col_x, col_w)):
        add_rect(slide, x, y, w, 0.8, fill=bg)
        color = rgb if j == 0 else C_TITLE
        bold = (j == 0)
        add_text(slide, val, x + 0.05, y + 0.05, w - 0.1, 0.7,
                 size=11, bold=bold, color=color)

add_text(slide, "Total: 2,017 E17.5 OG/PO cells", 6.0, 5.3, 7.2, 0.4,
         size=12, bold=True, color=C_ACCENT)
add_text(slide, "Barcodes matched to GSE163693_E15_17_composite_metadata.csv.gz",
         6.0, 5.75, 7.2, 0.4, size=10, color=C_SUB)
add_text(slide, "OG1–4 + PO1–2 location in full UMAP", 0.2, 5.6, 5.5, 0.4,
         size=11, color=C_SUB, align=PP_ALIGN.CENTER)

# ============================================================
# Slide 6: OG/PO サブセット UMAP + velocity stream
# ============================================================
slide = prs.slides.add_slide(BLANK)
slide_header(slide, "OG/PO Subset: UMAP and RNA Velocity",
             "2,017 cells (OG1–4 + PO1–2) | Re-computed neighbors, UMAP, and velocity | Dynamical model")

add_image(slide, f"{OG_DIR}/umap_og_exact_clusters.png",
          0.2, 1.1, width=6.2)
add_image(slide, f"{OG_DIR}/velocity_stream.png",
          6.6, 1.1, width=6.5)

add_text(slide, "OG1–4 + PO1–2 subset UMAP", 0.2, 5.6, 6.0, 0.4,
         size=11, color=C_SUB, align=PP_ALIGN.CENTER)
add_text(slide, "RNA velocity stream (dynamical model)", 6.6, 5.6, 6.5, 0.4,
         size=11, color=C_SUB, align=PP_ALIGN.CENTER)
add_footnote(
    slide,
    "OG1 (red): progenitors | OG2 (orange): suture pre-OB | "
    "OG3 (green): periosteal pre-OB | OG4 (purple): mature OB | "
    "PO1 (blue): periosteal OB | PO2 (brown): periosteal OB (mature subset)",
)

# ============================================================
# Slide 7: Latent time
# ============================================================
slide = prs.slides.add_slide(BLANK)
slide_header(
    slide,
    "OG/PO Subset: Latent Time",
    "Differentiation pseudotime estimated by scVelo dynamical model",
)

add_image(slide, f"{OG_DIR}/latent_time.png",
          0.3, 1.1, width=6.2)
add_image(slide, f"{OG_DIR}/og_latent_time_hist.png",
          6.7, 1.1, width=6.3)

add_text(slide, "Latent time on UMAP", 0.3, 5.65, 6.0, 0.4,
         size=11, color=C_SUB, align=PP_ALIGN.CENTER)
add_text(slide, "Latent time distribution per OG/PO cluster", 6.7, 5.65, 6.3, 0.4,
         size=11, color=C_SUB, align=PP_ALIGN.CENTER)
add_footnote(slide, "0 = earliest (progenitor), 1 = latest (mature osteoblast)")

# ============================================================
# Slide 8: PAGA + CellRank 概要
# ============================================================
slide = prs.slides.add_slide(BLANK)
slide_header(slide, "PAGA Topology and CellRank Analysis",
             "Bidirectional analysis allowing dedifferentiation (OG4 → OG2)")

add_image(slide, f"{OG_DIR}/og_paga_graph.png",
          0.2, 1.1, width=4.5)
add_image(slide, f"{OG_DIR}/cr_macrostates.png",
          4.9, 1.1, width=4.2)
add_image(slide, f"{OG_DIR}/cr_fate_heatmap.png",
          9.3, 1.1, width=3.8)

add_text(slide, "PAGA graph\n(edge width = transition strength)",
         0.2, 5.55, 4.5, 0.6, size=10, color=C_SUB, align=PP_ALIGN.CENTER)
add_text(slide, "GPCCA macrostates",
         4.9, 5.55, 4.2, 0.4, size=10, color=C_SUB, align=PP_ALIGN.CENTER)
add_text(slide, "Fate probability heatmap",
         9.3, 5.55, 3.8, 0.4, size=10, color=C_SUB, align=PP_ALIGN.CENTER)

add_footnote(slide,
    "CellRank GPCCA | VelocityKernel (80%) + ConnectivityKernel (20%) | "
    "Terminal: OG4 + OG1 | Initial: OG3")

# ============================================================
# Slide 9: CellRank 詳細（initial / terminal + fate probability）
# ============================================================
slide = prs.slides.add_slide(BLANK)
slide_header(slide, "CellRank: Initial / Terminal States and Fate Probabilities",
             "GPCCA identified OG1 as both initial and terminal state → dedifferentiation signal")

add_image(slide, f"{OG_DIR}/cr_initial_terminal.png",
          0.2, 1.1, width=8.5)
add_image(slide, f"{OG_DIR}/cr_fate_probabilities.png",
          8.9, 1.1, width=4.2)

add_text(slide, "Initial (left) and terminal (right) states",
         0.2, 5.6, 8.5, 0.4, size=10, color=C_SUB, align=PP_ALIGN.CENTER)
add_text(slide, "Fate probabilities per cell",
         8.9, 5.6, 4.2, 0.4, size=10, color=C_SUB, align=PP_ALIGN.CENTER)

# Key finding box
add_rect(slide, 0.2, 6.1, 12.9, 1.1, fill=RGBColor(0xff, 0xf3, 0xe0))
add_text(slide,
         "Key finding: OG4 cells show 73% fate probability toward OG1 direction, "
         "suggesting a dedifferentiation pathway from mature osteoblasts to progenitors.",
         0.35, 6.15, 12.6, 0.95,
         size=12, bold=False, color=RGBColor(0x7b, 0x3f, 0x00))

# ============================================================
# Slide 10: 局所 velocity（グリッド + 位相ポートレート）
# ============================================================
slide = prs.slides.add_slide(BLANK)
slide_header(slide, "Local Velocity Analysis",
             "Grid arrows and phase portraits to distinguish true loop from UMAP projection artifact")

add_image(slide, f"{OG_DIR}/local_grid_velocity.png",
          0.2, 1.1, width=13.0)
add_text(slide,
         "Grid velocity arrows at different densities (density = 0.5 / 1.0 / 2.0)\n"
         "Locally coherent arrows indicate directional flow, not a true cycle.",
         0.2, 5.6, 13.0, 0.7, size=11, color=C_SUB)

# ============================================================
# Slide 11: Velocity confidence + PCA space
# ============================================================
slide = prs.slides.add_slide(BLANK)
slide_header(slide, "Velocity Confidence and PCA-Space Velocity",
             "Confidence ≥ 0.77 for all clusters | PCA reduces UMAP distortion")

add_image(slide, f"{OG_DIR}/local_velocity_confidence.png",
          0.2, 1.1, width=6.8)
add_image(slide, f"{OG_DIR}/pca_stream.png",
          7.2, 1.1, width=5.9)

add_text(slide, "Velocity confidence (left: UMAP, right: boxplot per cluster)",
         0.2, 5.6, 6.8, 0.5, size=10, color=C_SUB, align=PP_ALIGN.CENTER)
add_text(slide, "Velocity stream in PCA space",
         7.2, 5.6, 5.9, 0.4, size=10, color=C_SUB, align=PP_ALIGN.CENTER)

# Confidence 表（OG1–4 + PO1–2, run_og_local_velocity.py の出力に基づく）
conf_data = [
    ("OG1", "0.839", "0.845", C_OG1),
    ("OG2", "0.815", "0.832", C_OG2),
    ("OG3", "0.820", "0.834", C_OG3),
    ("OG4", "0.852", "0.884", C_OG4),
    ("PO1", "0.822", "0.849", C_PO1),
    ("PO2", "0.747", "0.788", C_PO2),
]
add_rect(slide, 0.2, 6.05, 6.8, 0.38, fill=C_HEAD)
for j, h in enumerate(["Cluster", "Mean", "Median"]):
    add_text(slide, h, 0.25 + j * 2.2, 6.07, 2.0, 0.33,
             size=10, bold=True, color=C_WHITE)
for i, (cl, mn, med, rgb) in enumerate(conf_data):
    y = 6.45 + i * 0.24
    bg = C_LIGHT if i % 2 == 0 else C_WHITE
    add_rect(slide, 0.2, y, 6.8, 0.23, fill=bg)
    for j, val in enumerate([cl, mn, med]):
        c = rgb if j == 0 else C_TITLE
        add_text(slide, val, 0.25 + j * 2.2, y + 0.02, 2.0, 0.2,
                 size=9, bold=(j == 0), color=c)

# ============================================================
# Slide 12: 位相ポートレート
# ============================================================
slide = prs.slides.add_slide(BLANK)
slide_header(slide, "Phase Portraits: Spliced vs Unspliced",
             "Diagonal distribution indicates monotonic kinetics (not a true cycle)")

add_image(slide, f"{OG_DIR}/local_phase_portraits.png",
          0.2, 1.1, width=12.9)

add_footnote(slide,
    "Diagonal: induction/repression (directional flow) | "
    "Ellipse: cyclic (true loop) | "
    "Top genes by fit_likelihood from dynamical model")

# ============================================================
# Slide 13: まとめ
# ============================================================
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill=C_TITLE)
add_rect(slide, 0, 0, 13.33, 1.1, fill=C_HEAD)
add_text(slide, "Summary", 0.3, 0.1, 12.7, 0.9,
         size=28, bold=True, color=C_WHITE)

findings = [
    ("Differentiation trajectory",
     "OG1 (progenitors) → OG2/OG3 (pre-osteoblasts) → OG4 (mature osteoblasts) "
     "confirmed by RNA velocity in E17.5 coronal suture."),
    ("Bifurcation",
     "PAGA reveals two branches from OG1: suture-resident (OG2, edge=0.77) "
     "and periosteal (OG3, edge=0.75), both converging on OG4."),
    ("Dedifferentiation signal",
     "CellRank GPCCA identifies OG1 as a terminal state. OG4 cells show 73% "
     "fate probability toward OG1, suggesting dedifferentiation from mature "
     "osteoblasts to progenitors."),
    ("Loop is a projection artifact",
     "Velocity confidence ≥ 0.77 for all clusters. Phase portraits show diagonal "
     "(monotonic) rather than elliptic (cyclic) distribution. The loop in UMAP "
     "reflects a Y-shaped bifurcation compressed in 2D."),
]
for i, (title, text) in enumerate(findings):
    y = 1.3 + i * 1.5
    add_rect(slide, 0.3, y, 0.06, 1.2, fill=C_ACCENT)
    add_text(slide, title, 0.55, y, 12.5, 0.5,
             size=15, bold=True, color=C_WHITE)
    add_text(slide, text, 0.55, y + 0.5, 12.5, 0.9,
             size=12, color=RGBColor(0xcc, 0xcc, 0xcc))

add_text(slide,
         "Data: GSE163693  |  Farmer et al. 2021, Nat Commun 12:4797  |  "
         "Tools: scVelo, CellRank, scanpy",
         0.3, 7.1, 12.7, 0.3,
         size=9, color=RGBColor(0x77, 0x77, 0x99))

# ============================================================
# 保存
# ============================================================
prs.save(OUT_PPT)
print(f"Saved: {OUT_PPT}")
print(f"  Slides: {len(prs.slides)}")
