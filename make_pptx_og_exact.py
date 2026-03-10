"""
figures_og_exact 配下の図を順番に PowerPoint にまとめるスクリプト
"""

import os
from glob import glob

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = os.path.dirname(os.path.abspath(__file__))
os.chdir(WORKDIR)

FIG_DIR = "figures_og_exact"
OUT_PPT = "SutureRNAVelocity_OGExact_Figures.pptx"

# スライドサイズ：ワイド（16:9）
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

C_TITLE = RGBColor(0x1a, 0x1a, 0x2e)
C_HEAD = RGBColor(0x16, 0x21, 0x3e)
C_WHITE = RGBColor(0xff, 0xff, 0xff)
C_SUB = RGBColor(0x55, 0x55, 0x55)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

BLANK = prs.slide_layouts[6]


def add_rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(
        1, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text(
    slide,
    text,
    left,
    top,
    width,
    height,
    size=18,
    bold=False,
    color=C_TITLE,
    align=PP_ALIGN.LEFT,
):
    tx = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = tx.text_frame
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tx


def add_image_centered(slide, img_path):
    """画像をアスペクト比を保って、スライド中央に最大フィットで配置"""
    from PIL import Image

    img = Image.open(img_path)
    iw, ih = img.size  # px
    img.close()

    # スライド内のマージン
    margin_x = Inches(0.5)
    margin_y_top = Inches(1.2)
    margin_y_bottom = Inches(0.5)

    avail_w = SLIDE_W - 2 * margin_x
    avail_h = SLIDE_H - margin_y_top - margin_y_bottom

    img_ratio = iw / ih
    box_ratio = avail_w / avail_h

    if img_ratio > box_ratio:
        # 横長: 幅を合わせる
        width = avail_w
        height = width / img_ratio
    else:
        # 縦長: 高さを合わせる
        height = avail_h
        width = height * img_ratio

    left = (SLIDE_W - width) / 2
    top = margin_y_top + (avail_h - height) / 2

    slide.shapes.add_picture(img_path, left, top, width=width, height=height)


def slide_header(slide, title, subtitle=None):
    add_rect(slide, 0, 0, 13.33, 0.9, C_HEAD)
    add_text(
        slide,
        title,
        0.4,
        0.1,
        12.5,
        0.6,
        size=24,
        bold=True,
        color=C_WHITE,
        align=PP_ALIGN.LEFT,
    )
    if subtitle:
        add_text(
            slide,
            subtitle,
            0.4,
            0.55,
            12.5,
            0.4,
            size=12,
            color=RGBColor(0xcc, 0xcc, 0xcc),
            align=PP_ALIGN.LEFT,
        )


# 1枚目: タイトル
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, C_TITLE)
add_text(
    slide,
    "OG/PO Subset RNA Velocity Figures",
    0.5,
    2.8,
    12.3,
    1.0,
    size=32,
    bold=True,
    color=C_WHITE,
    align=PP_ALIGN.CENTER,
)
add_text(
    slide,
    "All figures from figures_og_exact/",
    0.5,
    3.9,
    12.3,
    0.6,
    size=16,
    color=RGBColor(0xdd, 0xdd, 0xee),
    align=PP_ALIGN.CENTER,
)
add_text(
    slide,
    "SutureRNAVelocity | GSE163693 | E17.5 OG1–4 + PO1–2",
    0.5,
    4.6,
    12.3,
    0.5,
    size=12,
    color=RGBColor(0xaa, 0xaa, 0xbb),
    align=PP_ALIGN.CENTER,
)

# 図ファイルのリストアップ（拡張子 png, jpg, pdf 等に対応）
patterns = ["*.png", "*.jpg", "*.jpeg", "*.tif", "*.tiff", "*.pdf"]
paths = []
for pat in patterns:
    paths.extend(glob(os.path.join(FIG_DIR, pat)))

# ファイル名でソート
paths = sorted(paths, key=lambda p: os.path.basename(p))

if not paths:
    print(f"No figure files found in {FIG_DIR}")
else:
    print(f"Found {len(paths)} figure files in {FIG_DIR}")

    for path in paths:
        fname = os.path.basename(path)
        slide = prs.slides.add_slide(BLANK)
        title = fname
        subtitle = None
        # 一部のよく使う図には簡単な説明を付与
        if "umap_og_exact_clusters" in fname:
            subtitle = "UMAP of OG1–4 + PO1–2 subset (cluster colors)"
        elif "velocity_stream" in fname and "local" not in fname:
            subtitle = "RNA velocity stream on OG/PO UMAP"
        elif "og_paga_graph" in fname:
            subtitle = "PAGA topology graph between OG/PO clusters"
        elif "cr_macrostates" in fname:
            subtitle = "CellRank GPCCA macrostates"
        elif "cr_diff_vector_field" in fname:
            subtitle = "Coarse-grained differentiation vector field (net flux)"
        elif "local_grid_velocity" in fname:
            subtitle = "Local grid velocity at different densities"
        elif "local_per_cluster_stream" in fname:
            subtitle = "Local velocity around each OG/PO cluster"
        elif "local_phase_portraits" in fname:
            subtitle = "Phase portraits (spliced vs unspliced) for top genes"
        elif "local_velocity_confidence" in fname:
            subtitle = "Velocity confidence per cell and per cluster"

        slide_header(slide, title, subtitle)
        add_image_centered(slide, path)

        # フッタ
        add_text(
            slide,
            f"{FIG_DIR}/{fname}",
            0.3,
            7.1,
            12.7,
            0.3,
            size=9,
            color=C_SUB,
            align=PP_ALIGN.RIGHT,
        )

prs.save(OUT_PPT)
print(f"Saved: {OUT_PPT}")
print(f"  Slides: {len(prs.slides)}")

